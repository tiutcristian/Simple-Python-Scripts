import os
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return "\n".join(text_content)


def write_text_to_pdf(text, pdf_path):
    c = canvas.Canvas(pdf_path, pagesize=letter)
    pdfmetrics.registerFont(TTFont('DejaVuSans', 'DejaVuSans.ttf'))
    c.setFont('DejaVuSans', 10)
    width, height = letter

    lines = text.split('\n')
    max_line_width = width - 80
    y = height - 40
    line_height = 14

    for line in lines:
        wrapped_lines = simpleSplit(line, 'DejaVuSans', 10, max_line_width)
        for wrapped_line in wrapped_lines:
            if y < 40:
                c.showPage()
                c.setFont('DejaVuSans', 10)
                y = height - 40
            c.drawString(40, y, wrapped_line)
            y -= line_height

    c.save()


def process_pptx_files_in_folder(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(folder_path, filename)
            pdf_path = os.path.join(folder_path, filename.replace(".pptx", ".pdf"))

            text = extract_text_from_pptx(pptx_path)
            write_text_to_pdf(text, pdf_path)

            print(f"Processed {pptx_path} to {pdf_path}")


folder_path = "C:/Users/tiutc/Desktop/CE"
process_pptx_files_in_folder(folder_path)
